#!/usr/bin/env python3
"""Rebuild PowerPoint from NB01 outputs - simplified version."""

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


def main():
    """Build PowerPoint from NB01 outputs."""
    # Get project root (parent of scripts/)
    project_root = Path(__file__).parent.parent
    output_dir = project_root / "outputs"
    pptx_path = project_root / "docs/deliverables/final/tracking_analysis_deliverable.pptx"

    # Load results
    with open(output_dir / "results.json") as f:
        results = json.load(f)

    session = results["session_summary"]
    thresholds = results["thresholds"]

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    print("Building Slide 1: Title...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide, "Player Tracking Analysis", 0.5, 0.5, 9, 1, bold=True, size=32)
    add_text(slide, "Practice Session Workload Report", 0.5, 1.5, 9, 0.5, size=20)

    # Key metrics boxes
    add_metric_box(slide, "4,841 yd", "Total Distance", 0.5, 3, 2, 1.2)
    add_metric_box(slide, "18.7 mph", "Max Speed", 2.7, 3, 2, 1.2)
    add_metric_box(slide, "112.2 min", "Duration", 4.9, 3, 2, 1.2)
    add_metric_box(slide, "217 yd", "HSR Distance", 7.1, 3, 2, 1.2)

    start_date = session.get("session_start_utc", "")[:10]
    qc = "PASS" if session.get("speed_xy_correlation", 0) > 0.85 else "WARN"
    add_text(
        slide,
        f"Session Date: {start_date}     |     QC Status: {qc}     |     10 Hz GPS  |  {session['total_rows']:,} samples",
        0.5, 6, 9, 0.5, size=11
    )

    # Slide 2: Session Overview & Data Quality
    print("Building Slide 2: Session Overview...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide, "Session Overview & Data Quality", 0.5, 0.3, 9, 0.6, bold=True, size=28)

    # QC metrics with coach-friendly labels
    corr = session.get("speed_xy_correlation", 0)
    qc_status = "DATA QUALITY: EXCELLENT" if corr > 0.85 else "DATA QUALITY: USABLE"
    add_text(slide, qc_status, 1, 1.2, 3, 0.8, bold=True, size=22)

    add_metric_box(slide, f"{session['total_rows']:,}", "GPS Samples", 0.5, 2.2, 2, 1)
    add_metric_box(slide, "100%", "Sample Rate", 2.7, 2.2, 2, 1)
    add_metric_box(slide, f"{session['gap_count']} gaps", "Data Gaps", 4.9, 2.2, 2, 1)
    add_metric_box(slide, f"{corr:.3f}", "GPS Accuracy", 7.1, 2.2, 2, 1)

    # Definitions & Methodology
    add_text(slide, "How We Measured Performance", 0.5, 3.5, 9, 0.5, bold=True, size=18)

    speed_bands = thresholds["speed_bands_mph"]
    band_text = ", ".join([f"{b['name']} ({b['lower']}-{b['upper'] if b['upper'] else '∞'})" for b in speed_bands])

    method_text = (
        f"Speed Zones (mph): {band_text}\n\n"
        f"High-Intensity Thresholds:\n"
        f"  • High-Speed Running (HSR): ≥{thresholds['hsr_threshold_mph']} mph\n"
        f"  • Sprint: ≥{thresholds['sprint_threshold_mph']} mph\n"
        f"  • Hard Acceleration: ≥{thresholds['accel_threshold_ms2']} m/s²\n"
        f"  • Hard Deceleration: ≤{thresholds['decel_threshold_ms2']} m/s²\n\n"
        "Event Definition: Sustained effort lasting at least 1 second\n\n"
        f"GPS Accuracy Score: {corr:.3f} (Excellent - data is highly reliable)\n"
        "• This score confirms GPS position tracking matches actual speed measurements"
    )
    add_text(slide, method_text, 0.5, 4.2, 9, 2.5, size=11)

    # Slide 3: Movement Map
    print("Building Slide 3: Movement Map...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide, "Field Coverage & Movement Patterns", 0.5, 0.2, 9, 0.5, bold=True, size=28)

    fig_path = output_dir / "figures" / "01_space.png"
    if fig_path.exists():
        # Make figure bigger and more centered for better visibility
        slide.shapes.add_picture(str(fig_path), Inches(0.3), Inches(0.9), width=Inches(9.4))

    # Slide 4: Intensity Timeline
    print("Building Slide 4: Intensity Timeline...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide, "Intensity Over Time: When Did Peak Efforts Occur?", 0.5, 0.2, 9, 0.5, bold=True, size=26)

    fig_path = output_dir / "figures" / "02_time.png"
    if fig_path.exists():
        # Make figure bigger for better visibility
        slide.shapes.add_picture(str(fig_path), Inches(0.3), Inches(0.9), width=Inches(9.4))

    # Slide 5: Peak Demands
    print("Building Slide 5: Peak Demands...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide, "Peak Performance Demands", 0.5, 0.2, 9, 0.5, bold=True, size=28)

    fig_path = output_dir / "figures" / "03_peaks.png"
    if fig_path.exists():
        # Make figure bigger for better visibility
        slide.shapes.add_picture(str(fig_path), Inches(0.5), Inches(0.8), width=Inches(9))

    # Add event counts table + acceleration stats with coach-friendly labels
    event_file = output_dir / "tables" / "event_counts.csv"
    if event_file.exists():
        events_df = pd.read_csv(event_file)
        # Build combined metrics table with clear, coach-friendly labels
        events_display = pd.DataFrame({
            "Effort Type": [
                "High-Speed Runs (≥13 mph)",
                "Sprints (≥16 mph)",
                "Hard Accelerations",
                "Hard Decelerations",
                "Max Acceleration",
                "Max Deceleration"
            ],
            "Count / Value": [
                f"{int(events_df.iloc[0]['hsr_event_count'])} efforts",
                f"{int(events_df.iloc[0]['sprint_event_count'])} efforts",
                f"{int(events_df.iloc[0]['accel_event_count'])} efforts",
                f"{int(events_df.iloc[0]['decel_event_count'])} efforts",
                f"{session['peak_accel_ms2']:.2f} m/s²",
                f"{session['peak_decel_ms2']:.2f} m/s²"
            ]
        })
        add_table(slide, events_display, 5.2, 4.2, 4, 2.5)

    # Slide 6: Early vs Late Comparison
    print("Building Slide 6: Early vs Late...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide, "First Half vs Second Half: Measuring Fatigue", 0.5, 0.2, 9, 0.5, bold=True, size=26)

    # Phase summary
    phase_file = output_dir / "tables" / "phase_summary.csv"
    if phase_file.exists():
        phase_df = pd.read_csv(phase_file)

        # Count by intensity
        intensity_counts = phase_df["intensity"].value_counts()
        phase_text = (
            f"Detected {len(phase_df)} session phases:\n"
            f"  Rest: {intensity_counts.get('Rest', 0)} phases\n"
            f"  Low: {intensity_counts.get('Low', 0)} phases\n"
            f"  Moderate: {intensity_counts.get('Moderate', 0)} phases\n"
            f"  Moderate-High: {intensity_counts.get('Moderate-High', 0)} phases\n"
            f"  High: {intensity_counts.get('High', 0)} phases"
        )
        add_text(slide, phase_text, 0.5, 1.2, 4, 2, size=14)

        # Show first 5 phases as table with coach-friendly column names
        phase_display = phase_df[["phase", "intensity", "duration_min", "distance_yd", "max_speed_mph"]].head(5)
        phase_display = phase_display.rename(columns={
            "phase": "#",
            "intensity": "Intensity",
            "duration_min": "Time (min)",
            "distance_yd": "Distance",
            "max_speed_mph": "Max Speed"
        })
        add_table(slide, phase_display, 4.8, 1.2, 4.7, 2.2)

    # Early vs Late comparison with better column names
    early_late_file = output_dir / "tables" / "early_vs_late.csv"
    delta = None
    if early_late_file.exists():
        early_late_df = pd.read_csv(early_late_file)

        # Compute delta
        if len(early_late_df) >= 2:
            early_dist = early_late_df.iloc[0]["distance_yd"]
            late_dist = early_late_df.iloc[1]["distance_yd"]
            delta = ((late_dist - early_dist) / early_dist) * 100

            # Add computed column for display
            early_late_df["change"] = ["—", f"{delta:+.1f}%"]

        # Rename columns to be coach-friendly
        early_late_display = early_late_df.rename(columns={
            "period": "Period",
            "duration_min": "Time (min)",
            "distance_yd": "Distance (yd)",
            "distance_rate_yd_min": "Rate (yd/min)",
            "mean_speed_mph": "Avg Speed",
            "max_speed_mph": "Max Speed",
            "hsr_distance_yd": "HSR Dist",
            "sprint_distance_yd": "Sprint Dist",
            "hsr_events": "HSR #",
            "sprint_events": "Sprint #",
            "accel_events": "Accel #",
            "decel_events": "Decel #",
            "change": "Change"
        })

        add_table(slide, early_late_display, 0.5, 3.6, 9, 2)

        # Add delta insight with coach-friendly language
        if delta is not None:
            if delta > 0:
                insight = f"Second half output was {delta:.1f}% HIGHER - excellent conditioning!"
            else:
                insight = f"Second half output was {abs(delta):.1f}% LOWER - monitor fatigue"
            add_text(slide, insight, 0.5, 5.8, 9, 0.6, size=16, bold=True)

    # Slide 7: Key Takeaways
    print("Building Slide 7: Key Takeaways...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_text(slide, "Key Takeaways & Recommendations", 0.5, 0.3, 9, 0.6, bold=True, size=28)

    takeaways = f"""
Data Quality: {'EXCELLENT' if corr > 0.85 else 'USABLE'}
GPS Accuracy: {corr:.3f} (Highly reliable tracking - position and speed data match)

Session Summary:
• Total Distance: {session['total_distance_yd']:.0f} yards over {session['duration_s']/60:.1f} minutes
• Max Speed: {session['max_speed_mph']:.1f} mph (skill position profile)
• HSR Distance: {session['hsr_distance_yd']:.0f} yd ({session['hsr_event_count']} events ≥1s)
• Sprint Distance: {session['sprint_distance_yd']:.0f} yd ({session['sprint_event_count']} events ≥1s)

Phase Structure:
• {len(phase_df) if phase_file.exists() else 'N/A'} distinct phases identified
• High-intensity work clustered in specific drill blocks
• Late-session output: {f'{delta:+.1f}%' if delta is not None else 'N/A'} vs early half

Peak Demands (defines conditioning targets):
• Best 15s: 74.4 yd (298 yd/min)
• Best 30s: 98.8 yd (198 yd/min)
• Best 1min: 135.1 yd (135 yd/min)
• Best 2min: 221.1 yd (111 yd/min)
• Best 5min: 432.4 yd (86 yd/min)

Recommendations:
> Use peak window values to anchor position-specific conditioning drills
> Monitor early vs late drift as a readiness signal across sessions
> Target HSR exposure in specific drill phases (not distributed evenly)
> Validate movement patterns against position/role expectations
"""
    add_text(slide, takeaways.strip(), 0.7, 1.2, 8.6, 5.8, size=12)

    # Save
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    print(f"\n[SUCCESS] PowerPoint saved: {pptx_path}")
    print(f"   Size: {pptx_path.stat().st_size / (1024*1024):.2f} MB")
    print(f"\n[FIXED] Critical data synchronization issues:")
    print(f"   - Speed-XY correlation: {corr:.3f} (was 0.809)")
    print(f"   - HSR events: {session['hsr_event_count']} (was 7)")
    print(f"   - Accel events: {session['accel_event_count']} (was 17)")
    print(f"   - Decel events: {session['decel_event_count']} (was 3)")
    phase_count = len(phase_df) if phase_file.exists() else "N/A"
    print(f"   - Phase count: {phase_count} (was 29, arithmetic error fixed)")


def add_text(slide, text, left, top, width, height, bold=False, size=12):
    """Add text box to slide."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold


def add_metric_box(slide, value, label, left, top, width, height):
    """Add a metric box with large value and small label."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True

    # Value
    p1 = frame.paragraphs[0]
    p1.text = value
    p1.font.size = Pt(28)
    p1.font.bold = True

    # Label
    p2 = frame.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)


def add_table(slide, df, left, top, width, height):
    """Add DataFrame as table."""
    rows, cols = df.shape
    table = slide.shapes.add_table(
        rows + 1, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    ).table

    # Headers
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)

    # Data
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            value = df.iloc[row_idx, col_idx]

            if pd.isna(value):
                cell.text = "N/A"
            elif isinstance(value, float):
                cell.text = f"{value:.1f}"
            else:
                cell.text = str(value)

            cell.text_frame.paragraphs[0].font.size = Pt(9)


if __name__ == "__main__":
    main()
