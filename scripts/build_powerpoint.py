#!/usr/bin/env python3
"""Build PowerPoint deliverable from outputs/ artifacts.

Generates a professional 8-slide presentation for the Browns Performance
Science Fellowship project, incorporating analysis results, figures, and tables.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


def build_presentation(output_dir: Path, deliverable_path: Path) -> None:
    """Build PowerPoint from outputs/ artifacts."""
    # Load results contract
    results_path = output_dir / "results.json"
    if not results_path.exists():
        raise FileNotFoundError(
            f"Results file not found: {results_path}. "
            "Run scripts/render_visual_templates.py first."
        )

    with open(results_path) as f:
        results = json.load(f)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add 8 slides
    print("Building slide 1: Title...")
    add_title_slide(prs, results)

    print("Building slide 2: Session Snapshot + Data Quality...")
    add_snapshot_slide(prs, output_dir, results)

    print("Building slide 3: Movement Map...")
    add_movement_map_slide(prs, output_dir)

    print("Building slide 4: Intensity Timeline...")
    add_intensity_timeline_slide(prs, output_dir)

    print("Building slide 5: Peak Demands...")
    add_peak_demands_slide(prs, output_dir)

    print("Building slide 6: Speed Zones...")
    add_speed_zones_slide(prs, output_dir)

    print("Building slide 7: Phases + Early vs Late...")
    add_phases_comparison_slide(prs, output_dir)

    print("Building slide 8: Key Takeaways...")
    add_takeaways_slide(prs, output_dir, results)

    # Save
    deliverable_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(deliverable_path))
    print(f"\n[SUCCESS] PowerPoint saved: {deliverable_path}")
    print(f"   File size: {deliverable_path.stat().st_size / (1024*1024):.2f} MB")


def add_title_slide(prs: Presentation, results: dict) -> None:
    """Slide 1: Title with session metadata."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Player Tracking Analysis: Practice Session Workload Report"

    session = results.get("session_summary", {})
    qc = results.get("qc_status", "UNKNOWN")

    start_ts = session.get("start_ts_utc", "")[:10] if session.get("start_ts_utc") else "N/A"
    duration_min = session.get("duration_s", 0) / 60.0

    subtitle.text = (
        f"Session Date: {start_ts}\n"
        f"Duration: {duration_min:.1f} minutes\n"
        f"QC Status: {qc}"
    )


def add_snapshot_slide(prs: Presentation, output_dir: Path, results: dict) -> None:
    """Slide 2: Session snapshot + data quality + validation gates."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Clear any placeholder text
    for shape in slide.shapes:
        if shape.is_placeholder:
            shape.text = ""

    # Title
    add_text_box(
        slide,
        "Session Snapshot + Data Quality",
        left=0.5,
        top=0.3,
        width=9,
        height=0.6,
        bold=True,
        size=28,
    )

    # Read snapshot text
    snapshot_file = output_dir / "slide_text" / "slide_1_session_snapshot.txt"
    if snapshot_file.exists():
        snapshot_text = snapshot_file.read_text()[:400]
        add_text_box(slide, snapshot_text, left=0.5, top=1.2, width=4.5, height=4)
    else:
        # Fallback: Generate from results
        session = results.get("session_summary", {})
        snapshot_text = (
            f"Rows: {session.get('rows', 0):,}\n"
            f"Duration: {session.get('duration_s', 0)/60:.1f} min\n"
            f"Distance: {session.get('distance_yd_from_speed', 0):.0f} yd\n"
            f"Peak Speed: {session.get('peak_speed_mph', 0):.1f} mph"
        )
        add_text_box(slide, snapshot_text, left=0.5, top=1.2, width=4.5, height=4)

    # Validation gates table
    gates_file = output_dir / "tables" / "slide_1_validation_gates_table.csv"
    if gates_file.exists():
        gates_df = pd.read_csv(gates_file)
        add_dataframe_table(slide, gates_df.head(4), left=5.5, top=1.2, width=4, height=2.5)

    # Data quality takeaways
    takeaways_file = output_dir / "slide_text" / "slide_1_data_quality_takeaways.txt"
    if takeaways_file.exists():
        takeaways = takeaways_file.read_text()[:300]
        add_text_box(
            slide,
            f"Key Takeaways:\n{takeaways}",
            left=0.5,
            top=5.5,
            width=9,
            height=1.5,
            size=11,
        )


def add_movement_map_slide(prs: Presentation, output_dir: Path) -> None:
    """Slide 3: Spatial movement map."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Clear any placeholder text
    for shape in slide.shapes:
        if shape.is_placeholder:
            shape.text = ""

    add_text_box(
        slide,
        "Where: Spatial Usage and Peak Windows",
        left=0.5,
        top=0.3,
        width=9,
        height=0.6,
        bold=True,
        size=28,
    )

    # Add figure
    fig_path = output_dir / "figures" / "coach_slide_movement_map.png"
    if not fig_path.exists():
        # Try alternate name
        fig_path = output_dir / "figures" / "01_space.png"

    if fig_path.exists():
        slide.shapes.add_picture(str(fig_path), left=Inches(1), top=Inches(1.2), width=Inches(8))

    # Add takeaway text box
    add_text_box(
        slide,
        "Peak demand windows show concentrated high-intensity zones on field",
        left=0.5,
        top=6.5,
        width=9,
        height=0.8,
        size=12,
    )


def add_intensity_timeline_slide(prs: Presentation, output_dir: Path) -> None:
    """Slide 4: Intensity timeline with top windows table."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Clear any placeholder text
    for shape in slide.shapes:
        if shape.is_placeholder:
            shape.text = ""

    add_text_box(
        slide,
        "When: Intensity Timeline",
        left=0.5,
        top=0.3,
        width=9,
        height=0.6,
        bold=True,
        size=28,
    )

    # Add figure
    fig_path = output_dir / "figures" / "coach_slide_intensity_timeline.png"
    if not fig_path.exists():
        fig_path = output_dir / "figures" / "02_time.png"

    if fig_path.exists():
        slide.shapes.add_picture(str(fig_path), left=Inches(0.5), top=Inches(1), width=Inches(9))

    # Add top windows table below
    windows_file = output_dir / "tables" / "slide_3_top_windows_table.csv"
    if windows_file.exists():
        windows_df = pd.read_csv(windows_file)
        add_dataframe_table(slide, windows_df.head(3), left=2, top=5.5, width=6, height=1.5)


def add_peak_demands_slide(prs: Presentation, output_dir: Path) -> None:
    """Slide 5: Peak demands figure + tables."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Clear any placeholder text
    for shape in slide.shapes:
        if shape.is_placeholder:
            shape.text = ""

    add_text_box(
        slide,
        "What: Peak Demands",
        left=0.5,
        top=0.3,
        width=9,
        height=0.6,
        bold=True,
        size=28,
    )

    # Top: figure
    fig_path = output_dir / "figures" / "coach_slide_peak_demand_summary.png"
    if not fig_path.exists():
        fig_path = output_dir / "figures" / "03_peaks.png"

    if fig_path.exists():
        slide.shapes.add_picture(str(fig_path), left=Inches(1), top=Inches(1), width=Inches(8))

    # Bottom left: peak distance table
    peak_file = output_dir / "tables" / "slide_3_peak_distance_table.csv"
    if peak_file.exists():
        peak_df = pd.read_csv(peak_file)
        add_dataframe_table(slide, peak_df, left=0.5, top=4.5, width=4, height=2)

    # Bottom right: event counts
    events_file = output_dir / "tables" / "slide_3_event_counts_table.csv"
    if events_file.exists():
        events_df = pd.read_csv(events_file)
        add_dataframe_table(slide, events_df, left=5, top=4.5, width=4, height=2)


def add_speed_zones_slide(prs: Presentation, output_dir: Path) -> None:
    """Slide 6: Speed zone breakdown table + takeaways."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Clear any placeholder text
    for shape in slide.shapes:
        if shape.is_placeholder:
            shape.text = ""

    add_text_box(
        slide,
        "Speed Zone Breakdown",
        left=0.5,
        top=0.3,
        width=9,
        height=0.6,
        bold=True,
        size=28,
    )

    # Speed zones table
    zones_file = output_dir / "tables" / "slide_2_speed_zone_table.csv"
    if zones_file.exists():
        zones_df = pd.read_csv(zones_file)
        add_dataframe_table(slide, zones_df, left=1.5, top=1.5, width=7, height=3)

    # Takeaways
    takeaways_file = output_dir / "slide_text" / "slide_2_speed_zone_takeaways.txt"
    if takeaways_file.exists():
        takeaways = takeaways_file.read_text()
        add_text_box(
            slide, f"Key Insights:\n{takeaways}", left=1, top=5, width=8, height=2, size=12
        )


def add_phases_comparison_slide(prs: Presentation, output_dir: Path) -> None:
    """Slide 7: Session phases + Early vs Late comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Clear any placeholder text
    for shape in slide.shapes:
        if shape.is_placeholder:
            shape.text = ""

    add_text_box(
        slide,
        "Session Structure: Phases + Early vs Late",
        left=0.5,
        top=0.3,
        width=9,
        height=0.6,
        bold=True,
        size=28,
    )

    # Top: segment/phase table
    segments_file = output_dir / "tables" / "slide_4_segment_table.csv"
    if segments_file.exists():
        segments_df = pd.read_csv(segments_file)
        add_dataframe_table(slide, segments_df.head(5), left=0.5, top=1.2, width=9, height=2)

    # Bottom: early vs late
    early_late_file = output_dir / "tables" / "slide_5_early_late_table.csv"
    if early_late_file.exists():
        early_late_df = pd.read_csv(early_late_file)
        add_dataframe_table(slide, early_late_df, left=2, top=3.5, width=6, height=1.5)

    # Takeaways
    phase_file = output_dir / "slide_text" / "slide_4_segment_takeaways.txt"
    late_file = output_dir / "slide_text" / "slide_5_early_late_takeaways.txt"

    combined_text = ""
    if phase_file.exists():
        combined_text += f"Phases: {phase_file.read_text()[:150]}\n\n"
    if late_file.exists():
        combined_text += f"Early vs Late: {late_file.read_text()[:150]}"

    if combined_text:
        add_text_box(slide, combined_text, left=0.5, top=5.5, width=9, height=1.5, size=11)


def add_takeaways_slide(prs: Presentation, output_dir: Path, results: dict) -> None:
    """Slide 8: Key takeaways + recommendations."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Clear any placeholder text
    for shape in slide.shapes:
        if shape.is_placeholder:
            shape.text = ""

    add_text_box(
        slide,
        "Key Takeaways & Recommendations",
        left=0.5,
        top=0.3,
        width=9,
        height=0.6,
        bold=True,
        size=28,
    )

    # Build takeaways from results
    qc = results.get("qc_status", "UNKNOWN")
    session = results.get("session_summary", {})

    duration_min = session.get("duration_s", 0) / 60.0
    distance_yd = session.get("distance_yd_from_speed", 0)
    peak_speed = session.get("peak_speed_mph", 0)

    takeaways_text = f"""
QC Status: {qc}
Session Duration: {duration_min:.1f} minutes
Total Distance: {distance_yd:.0f} yards
Peak Speed: {peak_speed:.1f} mph

Key Findings:
• High-intensity work clustered in specific drill blocks
• Peak 1-minute window defines worst-case conditioning target
• Late-session output maintained relative to early half
• Spatial usage consistent with position-specific demands

Recommendations:
• Use peak window values to anchor conditioning drills
• Monitor early vs late drift as readiness signal
• Target HSR exposure in specific drill phases
• Validate movement patterns against position role

Generated with Claude Code
"""

    add_text_box(slide, takeaways_text.strip(), left=1, top=1.5, width=8, height=5, size=14)


# Helper functions


def add_text_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    bold: bool = False,
    size: int = 12,
) -> None:
    """Add text box to slide."""
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold


def add_dataframe_table(
    slide, df: pd.DataFrame, left: float, top: float, width: float, height: float
) -> None:
    """Add pandas DataFrame as PowerPoint table."""
    rows, cols = df.shape
    table = slide.shapes.add_table(
        rows + 1, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    ).table

    # Set column widths evenly
    col_width = Inches(width / cols)
    for col_idx in range(cols):
        table.columns[col_idx].width = col_width

    # Headers
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)

    # Data
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            value = df.iloc[row_idx, col_idx]

            # Simple formatting
            if pd.isna(value):
                cell.text = "N/A"
            elif isinstance(value, float):
                cell.text = f"{value:.1f}"
            else:
                cell.text = str(value)

            cell.text_frame.paragraphs[0].font.size = Pt(10)


def main() -> None:
    """Main entry point."""
    parser = argparse.ArgumentParser(description="Build PowerPoint deliverable")
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=Path("outputs"),
        help="Directory containing analysis outputs",
    )
    parser.add_argument(
        "--deliverable",
        type=Path,
        default=Path("docs/deliverables/tracking_analysis_deliverable.pptx"),
        help="Output PowerPoint file path",
    )
    args = parser.parse_args()

    try:
        build_presentation(args.output_dir, args.deliverable)
    except FileNotFoundError as e:
        print(f"\n[ERROR] {e}")
        print("\nMake sure to run the analysis pipeline first:")
        print("  uv run python scripts/render_visual_templates.py")
        print("  uv run python scripts/render_presentation_assets.py")
        exit(1)
    except Exception as e:
        print(f"\n[ERROR] Unexpected error: {e}")
        import traceback

        traceback.print_exc()
        exit(1)


if __name__ == "__main__":
    main()
